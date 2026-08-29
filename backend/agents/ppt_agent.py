from pptx import Presentation

class PPTAgent:
    def analyze_pptx(self, file_path: str) -> dict:
        try:
            prs = Presentation(file_path)
            slides_info = []
            raw_text = []
            for slide in prs.slides:
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text
                slides_info.append({
                    "title": title,
                    "shapes": len(slide.shapes)
                })
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        raw_text.append(shape.text)
                        
            return {
                "type": "pptx",
                "slides_count": len(prs.slides),
                "slides": slides_info,
                "raw_text": "\n".join(raw_text)
            }
        except Exception as e:
            return {"error": str(e)}
