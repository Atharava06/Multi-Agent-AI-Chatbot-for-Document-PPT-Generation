from pptx import Presentation
from pptx.util import Inches
from services.gemini import gemini_service
import uuid
from pathlib import Path

class PPTGenerationAgent:
    def generate(self, context: str, template_filename: str = None) -> str:
        prompt = f"Generate ONLY the presentation slides based on this context:\n{context}\n\nCRITICAL: DO NOT generate any long-form document content, reports, or proposals. Extract ONLY the key points meant for a presentation.\n\nFormat output as strictly valid JSON without any markdown code blocks: [{{'title': 'Slide 1', 'content': 'bullet points'}}, ...]"
        content_json = gemini_service.generate_json(prompt)
        
        try:
            import json
            content_json = content_json.strip().strip("```json").strip("```").strip()
            slides_data = json.loads(content_json)
        except:
            slides_data = [{"title": "Generated Slide", "content": "Content generated"}]
            
        prs = Presentation()
        
        if template_filename:
            template_path = Path(f"uploads/{template_filename}")
            if template_path.exists():
                try:
                    prs = Presentation(template_path)
                except:
                    pass
                    
        import re
        for slide_data in slides_data:
            slide_layout = prs.slide_layouts[1] # title and content
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body = slide.shapes.placeholders[1]
            
            raw_title = slide_data.get('title', 'Slide')
            raw_content = slide_data.get('content', '')
            
            # Clean Markdown
            raw_title = re.sub(r'\*\*(.*?)\*\*', r'\1', raw_title).replace('#', '').strip()
            raw_content = re.sub(r'\*\*(.*?)\*\*', r'\1', raw_content).replace('#', '').strip()
            
            title.text = raw_title
            tf = body.text_frame
            tf.text = raw_content
            
        filename = f"generated_{uuid.uuid4().hex[:8]}.pptx"
        prs.save(f"generated/{filename}")
        
        return filename
